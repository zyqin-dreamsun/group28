import pptx
from pptx import Presentation
from typing import Dict, List, Any
import fitz  # PyMuPDF
from PIL import Image
import io
import base64
import asyncio
from concurrent.futures import ThreadPoolExecutor

class PPTParser:
    def __init__(self):
        self.executor = ThreadPoolExecutor(max_workers=4)
    
    async def parse_ppt(self, file_path: str) -> Dict[str, Any]:
        """解析PPT文件，提取结构和内容"""
        loop = asyncio.get_event_loop()
        
        if file_path.endswith('.pptx'):
            return await loop.run_in_executor(
                self.executor, 
                self._parse_pptx, 
                file_path
            )
        elif file_path.endswith('.pdf'):
            return await loop.run_in_executor(
                self.executor,
                self._parse_pdf,
                file_path
            )
        else:
            raise ValueError("Unsupported file format")
    
    def _parse_pptx(self, file_path: str) -> Dict[str, Any]:
        """解析PPTX文件"""
        prs = Presentation(file_path)
        pages = []
        
        for i, slide in enumerate(prs.slides):
            page_data = {
                "page_num": i + 1,
                "title": "",
                "text": "",
                "elements": [],
                "images": []
            }
            
            # 提取形状内容
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text.strip()
                    if text:
                        if shape.is_title or "Title" in shape.name:
                            page_data["title"] = text
                        else:
                            page_data["text"] += text + "\n"
                        
                        # 记录元素类型和位置
                        element = {
                            "type": "text",
                            "content": text,
                            "position": {
                                "left": shape.left,
                                "top": shape.top,
                                "width": shape.width,
                                "height": shape.height
                            }
                        }
                        page_data["elements"].append(element)
                
                # 提取图片
                if shape.shape_type == 13:  # 13表示图片
                    image_data = self._extract_image(shape)
                    if image_data:
                        page_data["images"].append(image_data)
            
            pages.append(page_data)
        
        # 提取目录信息
        toc = self._extract_toc(pages)
        
        return {
            "filename": file_path,
            "total_pages": len(pages),
            "toc": toc,
            "pages": pages
        }
    
    def _parse_pdf(self, file_path: str) -> Dict[str, Any]:
        """解析PDF文件（PPT另存为PDF的情况）"""
        doc = fitz.open(file_path)
        pages = []
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            
            # 提取文本
            text = page.get_text()
            
            # 提取图片
            images = []
            for img_index, img in enumerate(page.get_images(full=True)):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_data = {
                    "index": img_index,
                    "format": base_image["ext"],
                    "data": base64.b64encode(base_image["image"]).decode()
                }
                images.append(image_data)
            
            # 尝试提取标题（第一行或最大字体）
            lines = text.split('\n')
            title = lines[0] if lines else f"Page {page_num + 1}"
            
            page_data = {
                "page_num": page_num + 1,
                "title": title,
                "text": text,
                "images": images,
                "elements": []
            }
            
            pages.append(page_data)
        
        return {
            "filename": file_path,
            "total_pages": len(pages),
            "pages": pages
        }
    
    def _extract_image(self, shape) -> Dict:
        """提取PPT中的图片"""
        try:
            image = shape.image
            image_bytes = io.BytesIO(image.blob)
            img = Image.open(image_bytes)
            
            # 转换为base64
            buffered = io.BytesIO()
            img.save(buffered, format="PNG")
            img_str = base64.b64encode(buffered.getvalue()).decode()
            
            return {
                "format": "png",
                "data": img_str,
                "width": img.width,
                "height": img.height
            }
        except:
            return None
    
    def _extract_toc(self, pages: List[Dict]) -> List[Dict]:
        """从页面中提取目录结构"""
        toc = []
        for page in pages:
            if page["title"]:
                toc.append({
                    "page": page["page_num"],
                    "title": page["title"],
                    "level": self._determine_title_level(page["title"])
                })
        return toc
    
    def _determine_title_level(self, title: str) -> int:
        """判断标题层级"""
        if len(title) < 20 and title.isupper():
            return 1  # 主标题
        elif any(char.isdigit() for char in title) and '.' in title:
            return 2  # 编号标题
        else:
            return 3  # 普通标题